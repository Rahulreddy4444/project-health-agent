"""
Presentation Builder

Generates a 5-7 slide executive PowerPoint presentation for monthly reviews.
Uses python-pptx to create professional slides with RAG status visualization.
"""

import json
from datetime import datetime
from pathlib import Path
from typing import Dict, Any, List, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

import config


# ─── Color Palette ───────────────────────────────────────────────────────────

COLORS = {
    "bg_dark": RGBColor(0x1A, 0x1A, 0x2E),       # Deep navy
    "bg_card": RGBColor(0x16, 0x21, 0x3E),        # Card background
    "text_white": RGBColor(0xFF, 0xFF, 0xFF),
    "text_light": RGBColor(0xB0, 0xB0, 0xC0),
    "text_muted": RGBColor(0x80, 0x80, 0x90),
    "accent": RGBColor(0x00, 0xD2, 0xFF),          # Cyan accent
    "green": RGBColor(0x00, 0xC8, 0x53),
    "amber": RGBColor(0xFF, 0xB3, 0x00),
    "red": RGBColor(0xFF, 0x17, 0x44),
    "divider": RGBColor(0x30, 0x30, 0x50),
}

RAG_COLORS = {
    "Green": COLORS["green"],
    "Amber": COLORS["amber"],
    "Red": COLORS["red"],
    "Unknown": COLORS["text_muted"],
}


def _set_slide_bg(slide, color=None):
    """Set slide background to dark color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color or COLORS["bg_dark"]


def _add_text_box(slide, left, top, width, height, text, 
                  font_size=12, color=None, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a styled text box to a slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color or COLORS["text_white"]
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def _add_rag_circle(slide, left, top, size, rag_label):
    """Add a colored RAG indicator circle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top),
        Inches(size), Inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RAG_COLORS.get(rag_label, COLORS["text_muted"])
    shape.line.fill.background()
    return shape


def _add_rectangle(slide, left, top, width, height, fill_color=None, border_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


# ─── Slide Builders ──────────────────────────────────────────────────────────

def _build_title_slide(prs: Presentation, assessments: List[Dict]):
    """Slide 1: Title slide with month and portfolio snapshot."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    _set_slide_bg(slide)
    
    # Title
    _add_text_box(slide, 0.5, 1.0, 9.0, 1.0,
                  "Monthly Project Health Review",
                  font_size=32, bold=True, color=COLORS["text_white"],
                  alignment=PP_ALIGN.CENTER)
    
    # Date
    date_str = datetime.now().strftime("%B %Y")
    _add_text_box(slide, 0.5, 1.8, 9.0, 0.5,
                  date_str,
                  font_size=18, color=COLORS["accent"],
                  alignment=PP_ALIGN.CENTER)
    
    # Subtitle
    _add_text_box(slide, 0.5, 2.5, 9.0, 0.5,
                  "Professional Services — Executive Summary",
                  font_size=14, color=COLORS["text_light"],
                  alignment=PP_ALIGN.CENTER)
    
    # Portfolio snapshot at bottom
    total = len(assessments)
    green = sum(1 for a in assessments if a.get("overall_rag") == "Green")
    amber = sum(1 for a in assessments if a.get("overall_rag") == "Amber")
    red = sum(1 for a in assessments if a.get("overall_rag") == "Red")
    
    _add_text_box(slide, 0.5, 4.0, 9.0, 0.5,
                  f"Portfolio: {total} Active Projects",
                  font_size=16, color=COLORS["text_white"],
                  alignment=PP_ALIGN.CENTER)
    
    # RAG summary circles
    center_x = 3.5
    spacing = 1.5
    
    for i, (label, count, color) in enumerate([
        ("Green", green, COLORS["green"]),
        ("Amber", amber, COLORS["amber"]),
        ("Red", red, COLORS["red"]),
    ]):
        x = center_x + i * spacing
        _add_rag_circle(slide, x, 4.6, 0.5, label)
        _add_text_box(slide, x - 0.2, 5.2, 0.9, 0.3,
                      f"{count} {label}", font_size=11, color=color,
                      alignment=PP_ALIGN.CENTER)
    
    # Footer
    _add_text_box(slide, 0.5, 6.8, 9.0, 0.3,
                  "Confidential — For Internal Use Only",
                  font_size=9, color=COLORS["text_muted"],
                  alignment=PP_ALIGN.CENTER)


def _build_portfolio_overview_slide(prs: Presentation, assessments: List[Dict]):
    """Slide 2: Portfolio overview table with RAG status for all projects."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    
    _add_text_box(slide, 0.5, 0.3, 9.0, 0.6,
                  "Portfolio Overview",
                  font_size=24, bold=True, color=COLORS["text_white"])
    
    _add_text_box(slide, 0.5, 0.8, 9.0, 0.4,
                  "Health status across all active projects",
                  font_size=12, color=COLORS["text_light"])
    
    # Table
    rows = len(assessments) + 1
    cols = 7
    table_shape = slide.shapes.add_table(
        rows, cols,
        Inches(0.3), Inches(1.4),
        Inches(9.4), Inches(0.4 + len(assessments) * 0.55)
    )
    table = table_shape.table
    
    # Headers
    headers = ["Project", "PM", "Stage", "% Done", "Score", "RAG", "Confidence"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = COLORS["accent"]
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["bg_card"]
    
    # Column widths
    widths = [2.5, 1.3, 1.5, 0.7, 0.7, 0.7, 1.0]
    for i, w in enumerate(widths):
        table.columns[i].width = Inches(w)
    
    # Data rows
    for row_idx, a in enumerate(assessments, 1):
        snapshot = a.get("summary_snapshot", {})
        pct = snapshot.get("percent_complete")
        pct_str = f"{pct*100:.0f}%" if isinstance(pct, (int, float)) and pct <= 1 else str(pct or "N/A")
        
        values = [
            a.get("project_name", "Unknown")[:35],
            snapshot.get("project_manager", "N/A")[:18],
            snapshot.get("project_stage", "N/A")[:20],
            pct_str,
            f"{a.get('composite_score', 'N/A')}",
            a.get("overall_rag", "N/A"),
            a.get("confidence", "N/A"),
        ]
        
        for col_idx, val in enumerate(values):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9)
            
            # Color the RAG column
            if col_idx == 5:
                p.font.color.rgb = RAG_COLORS.get(val, COLORS["text_white"])
                p.font.bold = True
            else:
                p.font.color.rgb = COLORS["text_white"]
            
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["bg_dark"] if row_idx % 2 == 0 else COLORS["bg_card"]


def _build_trend_analysis_slide(prs: Presentation, assessments: List[Dict]):
    """Slide 3: Dimension-level health across projects."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    
    _add_text_box(slide, 0.5, 0.3, 9.0, 0.6,
                  "Health Dimensions — Cross-Project View",
                  font_size=24, bold=True, color=COLORS["text_white"])
    
    _add_text_box(slide, 0.5, 0.8, 9.0, 0.4,
                  "How each health dimension performs across the portfolio",
                  font_size=12, color=COLORS["text_light"])
    
    dim_names = ["Schedule Adherence", "Completion Progress", "Milestone Health",
                 "Task Risk Profile", "Stakeholder Signals"]
    
    # Create a grid: dimensions as rows, projects as columns
    y_start = 1.5
    row_height = 0.9
    
    for dim_idx, dim_name in enumerate(dim_names):
        y = y_start + dim_idx * row_height
        
        # Dimension label
        _add_text_box(slide, 0.3, y + 0.1, 2.5, 0.4,
                      dim_name, font_size=11, color=COLORS["text_white"], bold=True)
        
        # Project indicators
        for proj_idx, a in enumerate(assessments):
            x = 3.0 + proj_idx * 3.2
            
            # Find this dimension
            dim_data = None
            for d in a.get("dimensions", []):
                if d["dimension"] == dim_name:
                    dim_data = d
                    break
            
            if dim_data:
                label = dim_data.get("label", "Unknown")
                _add_rag_circle(slide, x, y + 0.1, 0.3, label)
                
                short_name = a.get("project_name", "")[:20]
                _add_text_box(slide, x + 0.4, y + 0.1, 2.5, 0.3,
                              f"{label}", font_size=9,
                              color=RAG_COLORS.get(label, COLORS["text_muted"]))
    
    # Project name labels at top
    for proj_idx, a in enumerate(assessments):
        x = 3.0 + proj_idx * 3.2
        short_name = a.get("project_name", "Unknown")[:25]
        _add_text_box(slide, x, 1.1, 3.0, 0.3,
                      short_name, font_size=10, color=COLORS["accent"], bold=True)
    
    # Summary insight
    dim_issues = {}
    for a in assessments:
        for d in a.get("dimensions", []):
            if d.get("label") in ("Red", "Amber"):
                dim_issues[d["dimension"]] = dim_issues.get(d["dimension"], 0) + 1
    
    if dim_issues:
        worst_dim = max(dim_issues.items(), key=lambda x: x[1])
        insight = f"⚠ Most common concern: {worst_dim[0]} (flagged in {worst_dim[1]} project{'s' if worst_dim[1] > 1 else ''})"
    else:
        insight = "✓ All dimensions are healthy across the portfolio"
    
    _add_rectangle(slide, 0.3, 6.2, 9.4, 0.5, fill_color=COLORS["bg_card"])
    _add_text_box(slide, 0.5, 6.25, 9.0, 0.4,
                  insight, font_size=11, color=COLORS["amber"] if dim_issues else COLORS["green"])


def _build_risk_spotlight_slide(prs: Presentation, assessments: List[Dict]):
    """Slide 4: Top emerging risks across the portfolio."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    
    _add_text_box(slide, 0.5, 0.3, 9.0, 0.6,
                  "Risk Spotlight",
                  font_size=24, bold=True, color=COLORS["text_white"])
    
    _add_text_box(slide, 0.5, 0.8, 9.0, 0.4,
                  "Top risks requiring leadership attention",
                  font_size=12, color=COLORS["text_light"])
    
    # Collect and rank all risks
    all_risks = []
    for a in assessments:
        for risk in a.get("key_risks", []):
            all_risks.append({
                "project": a.get("project_name", "Unknown"),
                "risk": risk,
                "overall_rag": a.get("overall_rag", "Unknown"),
            })
    
    # Also extract Red dimensions
    for a in assessments:
        for d in a.get("dimensions", []):
            if d.get("label") == "Red":
                all_risks.append({
                    "project": a.get("project_name", "Unknown"),
                    "risk": f"🔴 {d['dimension']}: {d.get('reasoning', 'See details')}",
                    "overall_rag": a.get("overall_rag", "Unknown"),
                })
    
    # Deduplicate by risk text
    seen = set()
    unique_risks = []
    for r in all_risks:
        if r["risk"] not in seen:
            seen.add(r["risk"])
            unique_risks.append(r)
    
    # Display top risks
    y = 1.5
    for i, risk in enumerate(unique_risks[:6]):
        _add_rectangle(slide, 0.3, y, 9.4, 0.7, fill_color=COLORS["bg_card"], border_color=COLORS["divider"])
        
        proj_name = risk["project"][:25]
        _add_text_box(slide, 0.5, y + 0.05, 2.0, 0.3,
                      proj_name, font_size=10, color=COLORS["accent"], bold=True)
        
        risk_text = risk["risk"][:120]
        _add_text_box(slide, 0.5, y + 0.3, 9.0, 0.35,
                      risk_text, font_size=9, color=COLORS["text_light"])
        
        y += 0.8
    
    if not unique_risks:
        _add_text_box(slide, 0.5, 2.5, 9.0, 0.5,
                      "✅ No critical risks identified across the portfolio.",
                      font_size=14, color=COLORS["green"],
                      alignment=PP_ALIGN.CENTER)


def _build_project_deep_dive_slide(prs: Presentation, assessment: Dict):
    """Slide 5-6: Deep dive for projects needing attention."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    
    rag = assessment.get("overall_rag", "Unknown")
    emoji_map = {"Green": "●", "Amber": "●", "Red": "●"}
    
    _add_text_box(slide, 0.5, 0.3, 8.0, 0.6,
                  f"Deep Dive: {assessment.get('project_name', 'Unknown')}",
                  font_size=22, bold=True, color=COLORS["text_white"])
    
    # RAG indicator
    _add_rag_circle(slide, 8.8, 0.35, 0.4, rag)
    
    snapshot = assessment.get("summary_snapshot", {})
    
    # Project info card
    _add_rectangle(slide, 0.3, 1.1, 4.2, 1.8, fill_color=COLORS["bg_card"])
    
    info_lines = [
        f"PM: {snapshot.get('project_manager', 'N/A')}",
        f"Stage: {snapshot.get('project_stage', 'N/A')}",
        f"Timeline: {snapshot.get('project_start', 'N/A')} → {snapshot.get('project_end', 'N/A')}",
        f"Completion: {snapshot.get('percent_complete', 'N/A')}",
        f"Composite Score: {assessment.get('composite_score', 'N/A')}/3.00",
    ]
    
    y = 1.2
    for line in info_lines:
        _add_text_box(slide, 0.5, y, 3.8, 0.3, line, font_size=10, color=COLORS["text_light"])
        y += 0.3
    
    # Dimension scores card
    _add_rectangle(slide, 4.8, 1.1, 5.0, 1.8, fill_color=COLORS["bg_card"])
    
    _add_text_box(slide, 5.0, 1.2, 4.5, 0.3,
                  "Dimension Scores", font_size=11, bold=True, color=COLORS["accent"])
    
    y = 1.5
    for d in assessment.get("dimensions", []):
        label = d.get("label", "Unknown")
        color = RAG_COLORS.get(label, COLORS["text_muted"])
        _add_rag_circle(slide, 5.0, y + 0.02, 0.2, label)
        _add_text_box(slide, 5.3, y, 4.3, 0.25,
                      f"{d['dimension']}: {label}",
                      font_size=9, color=color)
        y += 0.25
    
    # Key risks
    _add_text_box(slide, 0.5, 3.1, 9.0, 0.3,
                  "Key Risks", font_size=14, bold=True, color=COLORS["red"])
    
    y = 3.5
    for risk in assessment.get("key_risks", [])[:4]:
        _add_text_box(slide, 0.5, y, 9.0, 0.4,
                      f"• {risk[:100]}", font_size=9, color=COLORS["text_light"])
        y += 0.35
    
    # Recommendations
    _add_text_box(slide, 0.5, 5.0, 9.0, 0.3,
                  "Recommended Actions", font_size=14, bold=True, color=COLORS["accent"])
    
    y = 5.4
    for i, rec in enumerate(assessment.get("recommendations", [])[:4], 1):
        _add_text_box(slide, 0.5, y, 9.0, 0.4,
                      f"{i}. {rec[:100]}", font_size=9, color=COLORS["text_light"])
        y += 0.35


def _build_recommendations_slide(prs: Presentation, assessments: List[Dict], synthesis_text: str):
    """Slide 7: Portfolio-level recommendations and next steps."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    
    _add_text_box(slide, 0.5, 0.3, 9.0, 0.6,
                  "Recommendations & Next Steps",
                  font_size=24, bold=True, color=COLORS["text_white"])
    
    # Aggregate unique recommendations
    all_recs = []
    for a in assessments:
        all_recs.extend(a.get("recommendations", []))
    unique_recs = list(dict.fromkeys(all_recs))
    
    y = 1.2
    for i, rec in enumerate(unique_recs[:7], 1):
        _add_rectangle(slide, 0.3, y, 9.4, 0.5, fill_color=COLORS["bg_card"])
        _add_text_box(slide, 0.5, y + 0.08, 0.4, 0.3,
                      f"{i}.", font_size=12, color=COLORS["accent"], bold=True)
        _add_text_box(slide, 0.9, y + 0.08, 8.6, 0.35,
                      rec[:110], font_size=10, color=COLORS["text_white"])
        y += 0.6
    
    # Next steps footer
    _add_rectangle(slide, 0.3, 5.8, 9.4, 1.0, fill_color=COLORS["bg_card"], border_color=COLORS["accent"])
    _add_text_box(slide, 0.5, 5.9, 9.0, 0.3,
                  "Next Review", font_size=12, bold=True, color=COLORS["accent"])
    _add_text_box(slide, 0.5, 6.2, 9.0, 0.5,
                  "Weekly health reports will continue. Next monthly review scheduled in 4 weeks. "
                  "Please raise any concerns to the PM or steering committee immediately.",
                  font_size=10, color=COLORS["text_light"])


def build_monthly_presentation(assessments: List[Dict], synthesis_text: str = "",
                                output_path: Path = None) -> str:
    """
    Build the complete monthly executive presentation.
    
    Creates a 5-7 slide PowerPoint deck.
    """
    if output_path is None:
        output_path = config.MONTHLY_DIR / f"Monthly_Health_Review_{datetime.now().strftime('%Y_%m')}.pptx"
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    _build_title_slide(prs, assessments)
    
    # Slide 2: Portfolio Overview
    _build_portfolio_overview_slide(prs, assessments)
    
    # Slide 3: Trend / Dimension Analysis
    _build_trend_analysis_slide(prs, assessments)
    
    # Slide 4: Risk Spotlight
    _build_risk_spotlight_slide(prs, assessments)
    
    # Slides 5-6: Deep dives for non-Green projects (or all if ≤2 projects)
    attention_projects = [a for a in assessments if a.get("overall_rag") != "Green"]
    if not attention_projects:
        attention_projects = assessments  # Show all if everything is Green
    
    for a in attention_projects[:2]:  # Max 2 deep dives
        _build_project_deep_dive_slide(prs, a)
    
    # Slide 7: Recommendations
    _build_recommendations_slide(prs, assessments, synthesis_text)
    
    # Save
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    
    return str(output_path)


if __name__ == "__main__":
    # Test with sample data
    sample_assessments = [
        {
            "project_name": "Test Project A",
            "overall_rag": "Amber",
            "composite_score": 2.3,
            "confidence": "High",
            "dimensions": [
                {"dimension": "Schedule Adherence", "label": "Red", "weight": 0.3, "score": 1, "reasoning": "Test"},
                {"dimension": "Completion Progress", "label": "Amber", "weight": 0.25, "score": 2, "reasoning": "Test"},
                {"dimension": "Milestone Health", "label": "Green", "weight": 0.2, "score": 3, "reasoning": "Test"},
                {"dimension": "Task Risk Profile", "label": "Green", "weight": 0.15, "score": 3, "reasoning": "Test"},
                {"dimension": "Stakeholder Signals", "label": "Amber", "weight": 0.1, "score": 2, "reasoning": "Test"},
            ],
            "key_risks": ["Schedule delay risk"],
            "recommendations": ["Accelerate timeline"],
            "summary_snapshot": {
                "project_manager": "Test PM",
                "project_stage": "Build",
                "percent_complete": 0.5,
                "project_start": "2026-01-01",
                "project_end": "2026-12-31",
                "at_risk": "Medium",
            },
        }
    ]
    
    path = build_monthly_presentation(sample_assessments)
    print(f"✅ Test presentation saved to: {path}")
